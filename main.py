from pptx import Presentation
from pptx.util import Pt, Inches


def set_font_size(shape, font_size):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            font.size = Pt(font_size)


def create_ppt_from_paragraphs(paragraphs, font_size=24):
    prs = Presentation()
    # Use the "Blank" slide layout without a title or content placeholder
    slide_layout = prs.slide_layouts[6]

    for i, para in enumerate(paragraphs, start=1):
        # Check if the paragraph is not empty
        if para.strip():
            slide = prs.slides.add_slide(slide_layout)
            content = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(8), Inches(6))
            text_frame = content.text_frame
            text_frame.word_wrap = True

            p = text_frame.add_paragraph()
            p.text = para

            # Center align the text in the slide body
            p.alignment = 1  # 0=left, 1=center, 2=right

            # Set font size for the content
            set_font_size(content, font_size)

    return prs


if __name__ == "__main__":
    # Open the file "text.txt" in read mode and read its content
    with open("text.txt", "r") as file:
        text = file.read()

    paragraphs = text.split('\n\n')

    presentation = create_ppt_from_paragraphs(paragraphs, font_size=24)

    # Save the PowerPoint presentation as "output.pptx"
    presentation.save("output_del.pptx")
